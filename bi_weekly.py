import streamlit as st
import sqlite3
import os
import pandas as pd
import numpy as np
from datetime import datetime
import plotly.graph_objects as go
from datetime import timedelta
from pptx import Presentation
from pptx.util import Inches
import io
from pathlib import Path


downloads_dir = str(Path.home() / "Downloads")


st.title('Alamein Bi Weekly')


@st.cache_resource
def get_data():
    db_path = Path(__file__).parent.parent / "data" / "alamein_db.sqlite3"
    conn = sqlite3.connect(db_path)
    
    query_prod = ''' SELECT * FROM st_data '''
    query_total = ''' SELECT * FROM total_prod ORDER BY date'''
    query_tmp = ''' SELECT date,
losses AS "Production Differ (BOPD)",
remarks AS remarks 
FROM tmp_remarks
'''
    query_header = '''
SELECT unique_id, well_name, reservoir
FROM header_id
WHERE unique_id IN (
"Alamein-43:RS", "Alamein-42:BMB", "Alamein-44:AD",
"Alamein-48:KH", "Alamein-51:AD", "Alamein-52:KH",
"WAYID-1X:BMB", "NE-AL-25:RS", "AEB-3C-2X-3A:SD2&SD3",
"AEB-3C-5X:AD","Alamein-57:AD","Yidma-12:AD"
)
ORDER BY CASE unique_id
WHEN "Alamein-42:RS" THEN 1
WHEN "Alamein-43:BMB" THEN 2
WHEN "Alamein-44:AD" THEN 3
WHEN "Alamein-48:KH" THEN 4
WHEN "Alamein-51:AD" THEN 5
WHEN "Alamein-52:KH" THEN 6
WHEN "WAYID-1X:BMB" THEN 7
WHEN "NE-AL-25:RS" THEN 8
WHEN "AEB-3C-2X-3A:SD2&SD3" THEN 9
WHEN "AEB-3C-5X:AD" THEN 10
WHEN "Alamein-57:AD" THEN 11
WHEN "Yidma-12:AD" THEN 12
END
'''

    df_prod = pd.read_sql_query(query_prod, conn)
    df_total = pd.read_sql_query(query_total, conn)
    df_tmp = pd.read_sql_query(query_tmp, conn)
    df_header = pd.read_sql_query(query_header, conn)  # Added query for df_header
    conn.close()

    # Ensure the 'date' column is in datetime format
    df_prod['date'] = pd.to_datetime(df_prod['date'])
    df_total['date'] = pd.to_datetime(df_total['date'])
    df_tmp['date'] = pd.to_datetime(df_tmp['date'])
    df_prod=df_prod.sort_values('date')
    
    return df_prod, df_total, df_tmp, df_header

# Fetch data
df_prod, df_total, df_tmp, df_header = get_data()

# Extract min and max dates as datetime.date
min_date = df_total['date'].min().date()
max_date = df_total['date'].max().date()

col5, col6 = st.columns([1, 1])
with col5:
    use_slider = st.radio('Choose Date Input Method', ['Slider', 'Manual Input'])
    if use_slider == 'Slider':
        selected_date_range = st.slider(
            'Select date range',
            min_value=min_date,  # Min date as datetime.date
            max_value=max_date,  # Max date as datetime.date
            value=(min_date, max_date)  # Default range
        )
    else:
        # Manual input of dates
        start_date = st.date_input('Start Date', min_value=min_date, max_value=max_date, value=min_date)
        end_date = st.date_input('End Date', min_value=min_date, max_value=max_date, value=max_date)
        selected_date_range = (start_date, end_date)

# Define a function to filter the DataFrame based on the selected date range
@st.cache_resource
def apply_filters(df, selected_date_range):
    filtered_df = df[(df['date'] >= pd.to_datetime(selected_date_range[0])) & 
                        (df['date'] <= pd.to_datetime(selected_date_range[1]))]
    return filtered_df

# Apply the filter to df_prod
filtered_df_total = apply_filters(df_total, selected_date_range)
filtered_df_tmp = apply_filters(df_tmp, selected_date_range)
filtered_df_tmp['date'] = pd.to_datetime(filtered_df_tmp['date']).dt.strftime('%Y-%m-%d')


def create_production_chart(df_total):
    # Ensure the date column is in datetime format
    df_total['date'] = pd.to_datetime(df_total['date'])

    # Create the plotly figure
    fig = go.Figure()

    # Add the line and marker plot for net_oil
    fig.add_trace(go.Scatter(
        x=df_total['date'],
        y=df_total['net_oil'],
        mode='lines+markers',
        line=dict(color='green', width=2),
        marker=dict(size=8, color='green', symbol='square'),
        name='Net Oil'
    ))

    # Customize layout
    fig.update_layout(
        title={
            'text': "EHO Oil Production Performance",
            'y': 0.96,  # Title position
            'x': 0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': dict(size=29)
        },
        xaxis=dict(
            title='',
            tickformat='%d-%b-%Y',
            showgrid=True,
            gridcolor='lightgrey',
            linecolor='black',
            ticks='outside',
            tickangle=270,
            tickfont=dict(
                family='Arial',  # Optional: specify font family
                size=14,         # Adjust size if needed
                color='black',
                weight='bold'    # Ensure bold font
            )
        ),
        yaxis=dict(
            title="Oil Rate (BOPD)",
            tickfont=dict(
                family='Arial',  # Optional: specify font family
                size=14,         # Adjust size if needed
                color='black',
                weight='bold'    # Ensure bold font
            ),
            titlefont=dict(
                family='Arial',  # Optional: specify font family
                size=14,         # Adjust font size if needed
                color='black',
                weight='bold'    # Bold font
            ),
            showgrid=True,
            gridcolor='lightgrey',
            linecolor='black',
            ticks='outside',
            range=[4000, 6000],  # Adjust based on your data
            dtick=200  # Ensures 10 evenly spaced ticks between 4000 and 6000
        ),
        plot_bgcolor='white',
        margin=dict(l=50, r=50, t=80, b=80),
        height=600
    )

    # Show gridlines
    fig.update_xaxes(showgrid=True, gridcolor='lightgrey')
    fig.update_yaxes(showgrid=True, gridcolor='lightgrey')

    return fig


fig = create_production_chart(filtered_df_total)
st.plotly_chart(fig)

if 'remarks' in filtered_df_tmp.columns:
    filtered_df_tmp['remarks'] = filtered_df_tmp['remarks'].str.replace(r"\)", ")\n", regex=True)


filtered_df_tmp = filtered_df_tmp.reset_index(drop=True)

df_remakrs = st.data_editor(filtered_df_tmp,num_rows='dynamic',hide_index=True)


def create_unique_id_chart(filtered_df, unique_id, well_name, reservoir, start_date):
    # Ensure the 'date' column is in datetime format
    filtered_df['date'] = pd.to_datetime(filtered_df['date'])

    # Filter the data for the specific unique_id and start date
    unique_id_df = filtered_df[
        (filtered_df['unique_id'] == unique_id) &
        (filtered_df['date'] >= pd.to_datetime(start_date))
    ]
    if unique_id_df.empty:
        return None  # Skip if no data

    # Get the actual min and max dates from the filtered data
    min_date = unique_id_df['date'].min()
    max_date = unique_id_df['date'].max() + timedelta(days=10) # Use actual last date from the data

    # Calculate the range and interval for the primary y-axis (Net/Gross)
    max_primary_value = max(unique_id_df['net'].max(), unique_id_df['gross'].max())
    min_primary_value = 0  # Start from 0
    interval_primary = (max_primary_value - min_primary_value) / 10
    rounded_interval_primary = round(interval_primary, -1)  # Round to nearest 10
    yaxis_range_primary = [min_primary_value, rounded_interval_primary * 10]
    yaxis_ticks_primary = [i * rounded_interval_primary for i in range(11)]

    # Create a Plotly figure
    fig = go.Figure()

    # Add the "Net" line in green
    fig.add_trace(go.Scatter(
        x=unique_id_df['date'],
        y=unique_id_df['net'],
        mode='lines',  # Line only
        name='Net',
        line=dict(color='green', width=1.5),  # Thin line width
    ))

    # Add the "Gross" line in black
    fig.add_trace(go.Scatter(
        x=unique_id_df['date'],
        y=unique_id_df['gross'],
        mode='lines',  # Line only
        name='Gross',
        line=dict(color='black', width=1.5),  # Thin line width
    ))

    # Add the "WC" line on a secondary y-axis in blue
    fig.add_trace(go.Scatter(
        x=unique_id_df['date'],
        y=unique_id_df['wc'],
        mode='lines',  # Line only
        name='WC',
        line=dict(color='blue', width=1.5),  # Thin line width
        yaxis='y2'  # Secondary y-axis
    ))

    # Customize layout with secondary y-axis
    fig.update_layout(
        title={
            'text': f"<b>{well_name} Production Test & Performance</b><br><i>{reservoir}</i>",
            'y': 0.9,
            'x': 0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': dict(size=18)
        },
        xaxis=dict(
            showgrid=True,
            gridcolor='lightgrey',
            linecolor='black',
            tickfont=dict(
                family='Arial',
                size=16,
                color='black',
                weight='bold'
            ),
            titlefont=dict(
                family='Arial',
                size=16,
                color='black',
                weight='bold'
            ),
            ticks='outside',
            tickangle=270,  # Rotate labels vertically
            tickformat='%d-%b-%Y',  # Format: Day-Month-Year
            range=[min_date, max_date],  # Dynamically adjust range to actual data
            dtick=86400000.0 * 15  # Add a tick every 15 days (in milliseconds)
        ),
        yaxis=dict(
            title="Net/Gross (BBLS/day)",
            showgrid=True,
            gridcolor='lightgrey',
            linecolor='black',
            ticks='outside',
            titlefont=dict(
                family='Arial',
                size=16,
                color='black',
                weight='bold'
            ),
            tickfont=dict(
                family='Arial',
                size=16,
                color='black',
                weight='bold'
            ),
            tickmode='array',
            tickvals=yaxis_ticks_primary,
            range=yaxis_range_primary
        ),
        yaxis2=dict(
            title="WC (%)",
            overlaying='y',
            side='right',
            showgrid=False,
            ticks='outside',
            tickmode='array',
            tickvals=[i * 10 for i in range(11)],
            range=[0, 100],
            tickfont=dict(color='blue'),
            titlefont=dict(color='blue')
        ),
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.02,
            xanchor="right",
            x=1
        ),
        plot_bgcolor='white',
        margin=dict(l=50, r=50, t=120, b=80),
        height=600
    )

    # Add annotation for sum of net
    net_sum = int(unique_id_df['net'].sum())
    fig.add_annotation(
        x=0.01, y=1.15,
        xref="paper", yref="paper",
        text=f"<b>CUM: {net_sum:,} BBLS</b>",
        showarrow=False,
        font=dict(size=14, color="green", family="Arial"),
        align="left",
        bgcolor="white",
        bordercolor="green",
        borderwidth=2,
        borderpad=10
    )

    return fig


        
        


# Define the start date for each unique_id
start_dates = {
"Alamein-42:BMB": "2024-10-07",
"Alamein-43:RS": "2024-07-03",
"Alamein-44:AD": "2023-10-06",
"Alamein-48:KH": "2023-08-09",
"Alamein-51:AD": "2024-10-06",
"Alamein-52:KH": "2024-05-27",
"WAYID-1X:BMB": "2023-11-05",
"NE-AL-25:RS": "2023-02-22",
"AEB-3C-2X-3A:SD2&SD3": "2024-02-14",
"AEB-3C-5X:AD": "2024-07-10",
"Alamein-57:AD":"2025-02-11",
'Yidma-12:AD':'2025-04-08'
}


plots = []
for _, row in df_header.iterrows():
    unique_id = row['unique_id']
    well_name = row['well_name']
    reservoir = row['reservoir']
    start_date = start_dates.get(unique_id, "2024-01-01")
    
    fig = create_unique_id_chart(df_prod, unique_id, well_name, reservoir, start_date)
    if fig:
        st.plotly_chart(fig)
        plots.append((well_name, fig))


with  col6:
   if st.button("Make Presentation"):
    # Create a PowerPoint presentation
    presentation = Presentation()
    
    for well_name, fig in plots:
        # Adjust figure layout for better resolution
        fig.update_layout(
            title=dict(
                text=f"<b>{well_name} Production Test & Performance</b><br><i>{reservoir}</i>",  # Original title
                font=dict(size=24)
            ),
            xaxis=dict(
                title=dict(text="Date", font=dict(size=14)),  # Adjust x-axis title font size
                tickfont=dict(size=10)  # Adjust x-axis tick label font size
            ),
            yaxis=dict(
                title=dict(text="Net OIl & Gross bpd", font=dict(size=18)),  # Adjust y-axis title font size
                tickfont=dict(size=12)  # Adjust y-axis tick label font size
            ),
            legend=dict(font=dict(size=14)),
            width=1220,  # Set figure width in pixels
            height=620   # Set figure height in pixels
        )
        
        # Save each Plotly figure as an image in memory with high resolution
        image_stream = io.BytesIO()
        fig.write_image(image_stream, format="png", width=1200, height=800, scale=2)
        image_stream.seek(0)
        
        # Add a slide for each figure
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
        
        # Adjust positioning and sizing of the figure
        left = Inches(0.8)  # Distance from the left side of the slide
        top = Inches(0.2)   # Distance from the top of the slide
        width = Inches(10.0)  # Adjusted width to fit the slide better
        height = Inches(7.5)  # Adjusted height to fit the slide better

        slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
    
    # Save the PowerPoint to the Downloads directory
    ppt_file_path = os.path.join(downloads_dir, "Production_Presentation.pptx")
    presentation.save(ppt_file_path)
    
    st.success(f"Presentation saved to {ppt_file_path}")

    # Save the presentation to a BytesIO stream for downloading
    ppt_stream = io.BytesIO()
    presentation.save(ppt_stream)
    ppt_stream.seek(0)
    
    # Provide the PowerPoint as a downloadable file
    st.download_button(
        label="Download Presentation",
        data=ppt_stream,
        file_name="Production_Presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
